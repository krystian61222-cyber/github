#!/usr/bin/env python3
"""Generate Fintech presentation PPTX file with modern business design."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# === COLOR SCHEME ===
PRIMARY_DARK = RGBColor(0x1B, 0x28, 0x38)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_ORANGE = RGBColor(0xF4, 0xA2, 0x61)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
TEXT_DARK = RGBColor(0x2D, 0x34, 0x36)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_CARD = RGBColor(0x22, 0x33, 0x45)
LIGHT_TEAL = RGBColor(0xE0, 0xF7, 0xFA)
CARD_BORDER = RGBColor(0xDE, 0xE2, 0xE6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_transparency(shape, pct):
    """Set fill transparency (0=opaque, 100=transparent)."""
    from lxml import etree
    spPr = shape._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', str((100 - pct) * 1000))


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def circle(slide, l, t, sz, color, tr=70):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    set_transparency(s, tr)
    return s


def bar(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def badge(slide, num, l, t, sz=Inches(0.5), color=ACCENT_TEAL):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    return s


def card(slide, l, t, w, h, fill=TEXT_WHITE, brd=CARD_BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = brd
    s.line.width = Pt(1)
    return s


def txt(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, sz=Pt(16), bold=False, color=TEXT_DARK,
         align=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2), level=0):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = sz
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s1, PRIMARY_DARK)
circle(s1, Inches(-1.5), Inches(-1.5), Inches(4.5), ACCENT_TEAL, 88)
circle(s1, Inches(10.5), Inches(5.0), Inches(4.0), ACCENT_ORANGE, 88)
circle(s1, Inches(11.0), Inches(-2.0), Inches(3.0), ACCENT_TEAL, 92)
bar(s1, Inches(4.0), Inches(2.0), Inches(5.3), Inches(0.05), ACCENT_TEAL)

tf = txt(s1, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.5))
para(tf, u'\u201eEdukacja czy automatyzacja?\u201d', Pt(44), True, TEXT_WHITE, PP_ALIGN.CENTER)

tf2 = txt(s1, Inches(2.0), Inches(3.6), Inches(9.3), Inches(1.8))
para(tf2,
     u'Wp\u0142yw korzystania z aplikacji fintechowych i robo-doradc\u00f3w '
     u'na poziom wiedzy finansowej oraz \u015bwiadomo\u015b\u0107 ryzyka '
     u'inwestor\u00f3w indywidualnych.',
     Pt(18), False, RGBColor(0xAD, 0xD8, 0xE6), PP_ALIGN.CENTER)

bar(s1, Inches(4.0), Inches(5.4), Inches(5.3), Inches(0.05), ACCENT_ORANGE)

tf3 = txt(s1, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.8))
para(tf3, u'Projekt Badania Jako\u015bciowego: Fintech a \u015bwiadomo\u015b\u0107 inwestycyjna',
     Pt(13), False, RGBColor(0x99, 0xAA, 0xBB), PP_ALIGN.CENTER)

bar(s1, Inches(0.5), Inches(6.8), Inches(1.5), Inches(0.06), ACCENT_TEAL)


# ============================================================
# SLIDE 2: Uzasadnienie wyboru tematu (3-column cards)
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s2, LIGHT_BG)
bar(s2, Inches(0), Inches(0), Inches(0.1), SLIDE_H, ACCENT_TEAL)
bar(s2, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), PRIMARY_DARK)
circle(s2, Inches(11.0), Inches(-1.5), Inches(3.5), ACCENT_TEAL, 92)

tf = txt(s2, Inches(0.6), Inches(0.3), Inches(10.0), Inches(0.9))
para(tf, u'Uzasadnienie wyboru tematu', Pt(30), True, PRIMARY_DARK)
bar(s2, Inches(0.6), Inches(1.05), Inches(2.5), Inches(0.05), ACCENT_TEAL)

# Card 1
card(s2, Inches(0.5), Inches(1.4), Inches(3.9), Inches(5.6), TEXT_WHITE, CARD_BORDER)
bar(s2, Inches(0.5), Inches(1.4), Inches(3.9), Inches(0.08), ACCENT_TEAL)
tf = txt(s2, Inches(0.7), Inches(1.7), Inches(3.5), Inches(0.6))
para(tf, u'Dlaczego aktualny?', Pt(16), True, ACCENT_TEAL)
tf = txt(s2, Inches(0.7), Inches(2.3), Inches(3.5), Inches(4.5))
para(tf, u'\u2022 Boom na aplikacje inwestycyjne (Revolut, eToro, XTB) oraz robo-doradztwo (Finax, Portu)',
     Pt(13), False, TEXT_DARK, space_before=Pt(0))
para(tf, u'\u2022 Niskie bariery wej\u015bcia, brak minimalnych kwot, \u201egamifikacja\u201d interfejs\u00f3w',
     Pt(13), space_before=Pt(8))
para(tf, u'\u2022 Rodzi pytania o rzeczywist\u0105 \u015bwiadomo\u015b\u0107 ryzyka',
     Pt(13), space_before=Pt(8))

# Card 2
card(s2, Inches(4.7), Inches(1.4), Inches(3.9), Inches(5.6), TEXT_WHITE, CARD_BORDER)
bar(s2, Inches(4.7), Inches(1.4), Inches(3.9), Inches(0.08), ACCENT_ORANGE)
tf = txt(s2, Inches(4.9), Inches(1.7), Inches(3.5), Inches(0.6))
para(tf, u'Kogo dotyczy?', Pt(16), True, ACCENT_ORANGE)
tf = txt(s2, Inches(4.9), Inches(2.3), Inches(3.5), Inches(4.5))
para(tf, u'\u2022 Pocz\u0105tkuj\u0105cych inwestor\u00f3w detalicznych i gospodarstw domowych',
     Pt(13), False, TEXT_DARK, space_before=Pt(0))
para(tf, u'\u2022 Szukaj\u0105 alternatywy dla nisko oprocentowanych lokat',
     Pt(13), space_before=Pt(8))
para(tf, u'\u2022 Pierwszy kontakt z rynkiem kapita\u0142owym przez smartfon',
     Pt(13), space_before=Pt(8))

# Card 3
card(s2, Inches(8.9), Inches(1.4), Inches(3.9), Inches(5.6), TEXT_WHITE, CARD_BORDER)
bar(s2, Inches(8.9), Inches(1.4), Inches(3.9), Inches(0.08), PRIMARY_DARK)
tf = txt(s2, Inches(9.1), Inches(1.7), Inches(3.5), Inches(0.6))
para(tf, u'Jakie ma znaczenie?', Pt(16), True, PRIMARY_DARK)
tf = txt(s2, Inches(9.1), Inches(2.3), Inches(3.5), Inches(4.5))
para(tf, u'\u2022 Dla regulator\u00f3w (KNF) \u2013 skuteczno\u015b\u0107 ostrze\u017ce\u0144 i ochrona konsumenta',
     Pt(13), False, TEXT_DARK, space_before=Pt(0))
para(tf, u'\u2022 Dla rynku \u2013 ocena ryzyka \u201epaniki t\u0142umu\u201d',
     Pt(13), space_before=Pt(8))
para(tf, u'\u2022 Dla tw\u00f3rc\u00f3w aplikacji \u2013 odpowiedzialne projektowanie',
     Pt(13), space_before=Pt(8))


# ============================================================
# SLIDE 3: Cel badania i pytania badawcze (two-column)
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s3, TEXT_WHITE)
bar(s3, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)
bar(s3, Inches(0), Inches(0), Inches(0.1), SLIDE_H, PRIMARY_DARK)
circle(s3, Inches(-1.0), Inches(5.5), Inches(3.0), PRIMARY_DARK, 92)

tf = txt(s3, Inches(0.6), Inches(0.3), Inches(10.0), Inches(0.9))
para(tf, u'Cel badania i pytania badawcze', Pt(30), True, PRIMARY_DARK)
bar(s3, Inches(0.6), Inches(1.05), Inches(2.5), Inches(0.05), ACCENT_TEAL)

# Left column - Cel
card(s3, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.7), LIGHT_TEAL, ACCENT_TEAL)
tf = txt(s3, Inches(0.8), Inches(1.6), Inches(5.3), Inches(0.5))
para(tf, u'Cel:', Pt(18), True, ACCENT_TEAL, space_before=Pt(0))
tf = txt(s3, Inches(0.8), Inches(2.1), Inches(5.3), Inches(4.8))
para(tf,
     u'Zrozumienie, czy i w jaki spos\u00f3b interakcja z nowoczesnymi platformami '
     u'inwestycyjnymi kszta\u0142tuje rzeczywist\u0105 wiedz\u0119 finansow\u0105 u\u017cytkownik\u00f3w. '
     u'Zbadanie, czy aplikacje te pe\u0142ni\u0105 funkcj\u0119 edukacyjn\u0105, buduj\u0105c '
     u'\u015bwiadomo\u015b\u0107 rynkow\u0105, czy raczej zwalniaj\u0105 u\u017cytkownika z my\u015blenia, '
     u'usypiaj\u0105c czujno\u015b\u0107.',
     Pt(14), False, TEXT_DARK, space_before=Pt(0))

# Right column - Pytania
card(s3, Inches(6.6), Inches(1.4), Inches(6.2), Inches(5.7), TEXT_WHITE, CARD_BORDER)
tf = txt(s3, Inches(6.9), Inches(1.6), Inches(5.7), Inches(0.5))
para(tf, u'Pytania badawcze:', Pt(18), True, ACCENT_ORANGE, space_before=Pt(0))

badge(s3, '1', Inches(6.9), Inches(2.2), Inches(0.4), ACCENT_TEAL)
tf = txt(s3, Inches(7.5), Inches(2.15), Inches(5.0), Inches(1.3))
para(tf,
     u'W jaki spos\u00f3b interfejs i funkcjonalno\u015bci aplikacji finansowych wp\u0142ywaj\u0105 '
     u'na poczucie zrozumienia rynku i pewno\u015b\u0107 siebie inwestor\u00f3w indywidualnych?',
     Pt(12), False, TEXT_DARK, space_before=Pt(0))

badge(s3, '2', Inches(6.9), Inches(3.5), Inches(0.4), ACCENT_ORANGE)
tf = txt(s3, Inches(7.5), Inches(3.45), Inches(5.0), Inches(1.3))
para(tf,
     u'Dlaczego u\u017cytkownicy podejmuj\u0105 decyzje o zakupie konkretnych instrument\u00f3w \u2013 '
     u'w jakim stopniu opieraj\u0105 si\u0119 na w\u0142asnej wiedzy, a w jakim na sugestiach '
     u'algorytm\u00f3w i \u0142atwo\u015bci \u201eklikni\u0119cia\u201d?',
     Pt(12), False, TEXT_DARK, space_before=Pt(0))

badge(s3, '3', Inches(6.9), Inches(4.9), Inches(0.4), PRIMARY_DARK)
tf = txt(s3, Inches(7.5), Inches(4.85), Inches(5.0), Inches(1.5))
para(tf,
     u'Jak pocz\u0105tkuj\u0105cy inwestorzy rozumiej\u0105 i interpretuj\u0105 ryzyko utraty '
     u'kapita\u0142u w \u015brodowisku aplikacji, kt\u00f3re przypominaj\u0105 gry mobilne lub '
     u'platformy spo\u0142eczno\u015bciowe?',
     Pt(12), False, TEXT_DARK, space_before=Pt(0))


# ============================================================
# SLIDE 4: Charakterystyka badanych
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s4, LIGHT_BG)
bar(s4, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_ORANGE)
bar(s4, Inches(0), Inches(0), Inches(0.1), SLIDE_H, ACCENT_ORANGE)
circle(s4, Inches(10.8), Inches(5.0), Inches(3.5), ACCENT_ORANGE, 90)

tf = txt(s4, Inches(0.6), Inches(0.3), Inches(10.0), Inches(0.9))
para(tf, u'Charakterystyka badanych', Pt(30), True, PRIMARY_DARK)
bar(s4, Inches(0.6), Inches(1.05), Inches(2.5), Inches(0.05), ACCENT_ORANGE)

# Left card
card(s4, Inches(0.5), Inches(1.4), Inches(6.0), Inches(2.5), TEXT_WHITE, CARD_BORDER)
tf = txt(s4, Inches(0.8), Inches(1.55), Inches(5.5), Inches(2.2))
para(tf, u'Kto b\u0119dzie rozm\u00f3wc\u0105:', Pt(15), True, ACCENT_TEAL, space_before=Pt(0))
para(tf, u'\u2022 Osoby doros\u0142e (25\u201340 lat), zarz\u0105dzaj\u0105ce w\u0142asnym bud\u017cetem',
     Pt(13), space_before=Pt(6))
para(tf, u'\u2022 Bez wykszta\u0142cenia kierunkowego (finansowego/ekonomicznego)',
     Pt(13), space_before=Pt(4))

# Right card - Kryteria
card(s4, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.5), TEXT_WHITE, CARD_BORDER)
tf = txt(s4, Inches(7.1), Inches(1.55), Inches(5.5), Inches(2.2))
para(tf, u'Kryteria w\u0142\u0105czaj\u0105ce:', Pt(15), True, ACCENT_ORANGE, space_before=Pt(0))
para(tf, u'\u2022 Inwestowanie przez smartfon w ci\u0105gu ostatnich 12\u201324 miesi\u0119cy',
     Pt(13), space_before=Pt(6))
para(tf, u'\u2022 Minimum 5 zrealizowanych transakcji',
     Pt(13), space_before=Pt(4))

# Bottom left card
card(s4, Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.8), TEXT_WHITE, CARD_BORDER)
tf = txt(s4, Inches(0.8), Inches(4.35), Inches(5.5), Inches(2.5))
para(tf, u'Wielko\u015b\u0107 pr\u00f3by:', Pt(15), True, PRIMARY_DARK, space_before=Pt(0))
para(tf, u'\u2022 12\u201315 os\u00f3b (nasycenie informacyjne)', Pt(13), space_before=Pt(6))
para(tf, u'Dlaczego ta grupa:', Pt(15), True, PRIMARY_DARK, space_before=Pt(12))
para(tf, u'\u2022 \u201e\u015awie\u017cy\u201d inwestorzy najbardziej nara\u017ceni na skutki braku edukacji',
     Pt(13), space_before=Pt(6))

# Bottom right card
card(s4, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.8), TEXT_WHITE, CARD_BORDER)
tf = txt(s4, Inches(7.1), Inches(4.35), Inches(5.5), Inches(2.5))
para(tf, u'Dob\u00f3r pr\u00f3by:', Pt(15), True, PRIMARY_DARK, space_before=Pt(0))
para(tf, u'\u2022 Celowy (selekcja wg kryteri\u00f3w sta\u017cu i aplikacji)',
     Pt(13), space_before=Pt(6))
para(tf, u'\u2022 Po\u0142\u0105czony z metod\u0105 kuli \u015bnie\u017cnej',
     Pt(13), space_before=Pt(4))


# ============================================================
# SLIDE 5: Dlaczego IDI? (dark background variant)
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s5, PRIMARY_DARK)
circle(s5, Inches(10.0), Inches(-1.5), Inches(4.0), ACCENT_TEAL, 90)
circle(s5, Inches(-1.0), Inches(5.5), Inches(3.5), ACCENT_ORANGE, 92)
bar(s5, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_TEAL)

tf = txt(s5, Inches(0.8), Inches(0.4), Inches(11.0), Inches(0.9))
para(tf, u'Dlaczego wywiad pog\u0142\u0119biony (IDI)?', Pt(30), True, TEXT_WHITE)
bar(s5, Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.05), ACCENT_TEAL)

reasons = [
    u'Finanse osobiste (straty, b\u0142\u0119dy, brak wiedzy) to tematy mocno wra\u017cliwe \u2013 '
    u'w grupie uczestnicy ulegaliby presji spo\u0142ecznej',
    u'Zale\u017cy nam na bardzo dok\u0142adnym prze\u015bledzeniu zachowania u\u017cytkownika '
    u'i jego interakcji z aplikacj\u0105',
    u'Chcemy, \u017ceby rozm\u00f3wca odtworzy\u0142 proces podejmowania konkretnej decyzji '
    u'inwestycyjnej krok po kroku',
    u'Tylko w IDI mamy swobod\u0119, \u017ceby na bie\u017c\u0105co reagowa\u0107, dopytywa\u0107 '
    u'o detale i pog\u0142\u0119bia\u0107 w\u0105tki',
]

for i, reason in enumerate(reasons):
    top = Inches(1.5 + i * 1.4)
    card(s5, Inches(0.6), top, Inches(12.0), Inches(1.15), DARK_CARD, DARK_CARD)
    badge(s5, str(i + 1), Inches(0.9), top + Inches(0.3), Inches(0.45), ACCENT_TEAL)
    tf = txt(s5, Inches(1.6), top + Inches(0.15), Inches(10.5), Inches(1.0))
    para(tf, reason, Pt(15), False, TEXT_WHITE, space_before=Pt(0))


# ============================================================
# SLIDE 6: Scenariusz wywiadu - A. Wprowadzenie
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s6, TEXT_WHITE)
bar(s6, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)
bar(s6, Inches(0), Inches(0), Inches(0.1), SLIDE_H, PRIMARY_DARK)
circle(s6, Inches(11.0), Inches(5.0), Inches(3.0), ACCENT_TEAL, 92)

tf = txt(s6, Inches(0.6), Inches(0.3), Inches(11.0), Inches(0.9))
para(tf, u'Scenariusz wywiadu \u2013 A. Wprowadzenie', Pt(28), True, PRIMARY_DARK)
bar(s6, Inches(0.6), Inches(1.0), Inches(2.5), Inches(0.05), ACCENT_TEAL)

# Section label
card(s6, Inches(0.5), Inches(1.3), Inches(2.0), Inches(0.5), ACCENT_TEAL, ACCENT_TEAL)
tf = txt(s6, Inches(0.6), Inches(1.33), Inches(1.8), Inches(0.45))
para(tf, u'CZ\u0118\u015a\u0106 A', Pt(11), True, TEXT_WHITE, PP_ALIGN.CENTER, space_before=Pt(0))

q6 = [
    u'Jakie by\u0142y Pana/Pani g\u0142\u00f3wne motywacje, \u017ceby w og\u00f3le zacz\u0105\u0107 '
    u'lokowa\u0107 swoje oszcz\u0119dno\u015bci poza standardowym kontem w banku?',
    u'Z jakich aplikacji lub platform inwestycyjnych Pan(i) obecnie korzysta '
    u'i dlaczego pad\u0142 wyb\u00f3r akurat na nie?',
]

for i, q in enumerate(q6):
    top = Inches(2.1 + i * 2.2)
    card(s6, Inches(0.5), top, Inches(12.0), Inches(1.8), LIGHT_BG, CARD_BORDER)
    badge(s6, str(i + 1), Inches(0.8), top + Inches(0.6), Inches(0.55), ACCENT_TEAL)
    tf = txt(s6, Inches(1.6), top + Inches(0.3), Inches(10.5), Inches(1.4))
    para(tf, q, Pt(16), False, TEXT_DARK, space_before=Pt(0))


# ============================================================
# SLIDE 7: Scenariusz wywiadu - B. Czesc glowna
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s7, LIGHT_BG)
bar(s7, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_ORANGE)
bar(s7, Inches(0), Inches(0), Inches(0.1), SLIDE_H, ACCENT_ORANGE)
circle(s7, Inches(-1.0), Inches(-1.0), Inches(3.0), ACCENT_ORANGE, 92)

tf = txt(s7, Inches(0.6), Inches(0.2), Inches(11.0), Inches(0.8))
para(tf, u'Scenariusz wywiadu \u2013 B. Cz\u0119\u015b\u0107 g\u0142\u00f3wna', Pt(26), True, PRIMARY_DARK)
bar(s7, Inches(0.6), Inches(0.85), Inches(2.5), Inches(0.04), ACCENT_ORANGE)

card(s7, Inches(0.5), Inches(1.1), Inches(2.0), Inches(0.45), ACCENT_ORANGE, ACCENT_ORANGE)
tf = txt(s7, Inches(0.6), Inches(1.12), Inches(1.8), Inches(0.4))
para(tf, u'CZ\u0118\u015a\u0106 B', Pt(10), True, TEXT_WHITE, PP_ALIGN.CENTER, space_before=Pt(0))

q7 = [
    u'Prosz\u0119 opowiedzie\u0107 mi o swojej pierwszej transakcji lub za\u0142o\u017ceniu portfela '
    u'w aplikacji. Jak to wygl\u0105da\u0142o krok po kroku i jak si\u0119 Pan(i) wtedy czu\u0142(a)?',
    u'W jaki spos\u00f3b aplikacja, z kt\u00f3rej Pan(i) korzysta, pomaga (lub nie) zrozumie\u0107, '
    u'w co dok\u0142adnie inwestowane s\u0105 pieni\u0105dze?',
    u'Zdarzy\u0142o si\u0119 Panu(i) podj\u0105\u0107 decyzj\u0119 o zainwestowaniu w co\u015b pod wp\u0142ywem impulsu, '
    u'bo zobaczy\u0142(a) Pan(i) to na ekranie g\u0142\u00f3wnym aplikacji?',
    u'Gdyby mia\u0142(a) Pan(i) swoimi s\u0142owami wyja\u015bni\u0107, czym jest instrument, w kt\u00f3ry '
    u'Pan(i) inwestuje (np. ETF, akcja) \u2013 jak by to Pan(i) opisa\u0142(a)?',
    u'Aplikacje cz\u0119sto pokazuj\u0105 komunikaty o ryzyku utraty kapita\u0142u. Jak Pan(i) reaguje '
    u'na te komunikaty w momencie podejmowania decyzji o zakupie?',
    u'Co si\u0119 dzieje, gdy wchodzi Pan(i) do aplikacji i widzi, \u017ce ca\u0142y wykres portfela jest '
    u'\u201ena czerwono\u201d? Co Pan(i) wtedy robi i dlaczego?',
    u'Czy uwa\u017ca Pan(i), \u017ce od kiedy u\u017cywa Pan(i) tej aplikacji, Pana/Pani wiedza '
    u'o finansach i ekonomii wzros\u0142a? Sk\u0105d to przekonanie?',
]

for i, q in enumerate(q7):
    top = Inches(1.65 + i * 0.82)
    badge(s7, str(i + 3), Inches(0.6), top + Inches(0.05), Inches(0.4), ACCENT_ORANGE)
    tf = txt(s7, Inches(1.2), top, Inches(11.5), Inches(0.8))
    para(tf, q, Pt(12), False, TEXT_DARK, space_before=Pt(0))


# ============================================================
# SLIDE 8: Scenariusz wywiadu - C. Zakonczenie
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s8, TEXT_WHITE)
bar(s8, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY_DARK)
bar(s8, Inches(0), Inches(0), Inches(0.1), SLIDE_H, PRIMARY_DARK)
circle(s8, Inches(10.5), Inches(-1.5), Inches(3.5), PRIMARY_DARK, 92)

tf = txt(s8, Inches(0.6), Inches(0.3), Inches(11.0), Inches(0.9))
para(tf, u'Scenariusz wywiadu \u2013 C. Zako\u0144czenie', Pt(28), True, PRIMARY_DARK)
bar(s8, Inches(0.6), Inches(1.0), Inches(2.5), Inches(0.05), PRIMARY_DARK)

card(s8, Inches(0.5), Inches(1.3), Inches(2.0), Inches(0.5), PRIMARY_DARK, PRIMARY_DARK)
tf = txt(s8, Inches(0.6), Inches(1.33), Inches(1.8), Inches(0.45))
para(tf, u'CZ\u0118\u015a\u0106 C', Pt(11), True, TEXT_WHITE, PP_ALIGN.CENTER, space_before=Pt(0))

q8 = [
    u'Czy uwa\u017ca Pan(i), \u017ce takie aplikacje powinny przed pozwoleniem na zakup '
    u'sprawdza\u0107 wiedz\u0119 u\u017cytkownika (np. kr\u00f3tkim quizem), czy inwestowanie powinno '
    u'by\u0107 dost\u0119pne dla ka\u017cdego bez ogranicze\u0144? Dlaczego?',
    u'Gdyby znajomy bez \u017cadnej wiedzy o finansach zapyta\u0142 Pana(i\u0105), czy warto '
    u'zacz\u0105\u0107 korzysta\u0107 z tej samej aplikacji \u2013 na co uwa\u017ca\u0142by/aby Pan(i) '
    u'za konieczne, by go ostrzec?',
    u'Zbli\u017camy si\u0119 do ko\u0144ca naszej rozmowy. Czy jest jeszcze co\u015b wa\u017cnego '
    u'w kwestii Pana/Pani do\u015bwiadcze\u0144 z tymi aplikacjami, o co nie zapyta\u0142em?',
]

for i, q in enumerate(q8):
    top = Inches(2.0 + i * 1.7)
    card(s8, Inches(0.5), top, Inches(12.0), Inches(1.4), LIGHT_BG, CARD_BORDER)
    badge(s8, str(i + 10), Inches(0.8), top + Inches(0.4), Inches(0.55), PRIMARY_DARK)
    tf = txt(s8, Inches(1.6), top + Inches(0.2), Inches(10.5), Inches(1.2))
    para(tf, q, Pt(15), False, TEXT_DARK, space_before=Pt(0))


# ============================================================
# SLIDE 9: Ograniczenia badania (3 cards)
# ============================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s9, PRIMARY_DARK)
circle(s9, Inches(-1.0), Inches(-1.0), Inches(4.0), ACCENT_TEAL, 90)
circle(s9, Inches(11.0), Inches(5.0), Inches(3.5), ACCENT_ORANGE, 90)
bar(s9, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_ORANGE)

tf = txt(s9, Inches(0.8), Inches(0.4), Inches(11.0), Inches(0.9))
para(tf, u'Ograniczenia badania', Pt(30), True, TEXT_WHITE)
bar(s9, Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.05), ACCENT_ORANGE)

# Card 1 - Dunning-Kruger
card(s9, Inches(0.4), Inches(1.5), Inches(4.0), Inches(5.3), DARK_CARD, DARK_CARD)
bar(s9, Inches(0.4), Inches(1.5), Inches(4.0), Inches(0.08), ACCENT_TEAL)
tf = txt(s9, Inches(0.7), Inches(1.8), Inches(3.5), Inches(0.5))
para(tf, u'Efekt Dunninga-Krugera', Pt(14), True, ACCENT_TEAL, space_before=Pt(0))
tf = txt(s9, Inches(0.7), Inches(2.4), Inches(3.5), Inches(4.0))
para(tf, u'\u2022 Badani mog\u0105 znacz\u0105co przecenia\u0107 swoj\u0105 wiedz\u0119 finansow\u0105, '
     u'szczeg\u00f3lnie w okresach hossy',
     Pt(12), False, TEXT_WHITE, space_before=Pt(0))
para(tf, u'\u2022 Deklaratywny wzrost wiedzy mo\u017ce by\u0107 w rzeczywisto\u015bci tylko '
     u'z\u0142udzeniem kontroli',
     Pt(12), False, TEXT_WHITE, space_before=Pt(8))

# Card 2 - Trudnosc ewaluacji
card(s9, Inches(4.7), Inches(1.5), Inches(4.0), Inches(5.3), DARK_CARD, DARK_CARD)
bar(s9, Inches(4.7), Inches(1.5), Inches(4.0), Inches(0.08), ACCENT_ORANGE)
tf = txt(s9, Inches(5.0), Inches(1.8), Inches(3.5), Inches(0.5))
para(tf, u'Trudno\u015b\u0107 ewaluacji wiedzy', Pt(14), True, ACCENT_ORANGE, space_before=Pt(0))
tf = txt(s9, Inches(5.0), Inches(2.4), Inches(3.5), Inches(4.0))
para(tf, u'\u2022 Wywiad jako\u015bciowy nie jest testem wiedzy',
     Pt(12), False, TEXT_WHITE, space_before=Pt(0))
para(tf, u'\u2022 Trudno wywa\u017cy\u0107 dopytywanie o definicje, tak aby badany '
     u'nie poczu\u0142 si\u0119 jak na egzaminie',
     Pt(12), False, TEXT_WHITE, space_before=Pt(8))

# Card 3 - Bariera technologiczna
card(s9, Inches(9.0), Inches(1.5), Inches(4.0), Inches(5.3), DARK_CARD, DARK_CARD)
bar(s9, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.08), TEXT_WHITE)
tf = txt(s9, Inches(9.3), Inches(1.8), Inches(3.5), Inches(0.5))
para(tf, u'Bariera technologiczna', Pt(14), True, TEXT_WHITE, space_before=Pt(0))
tf = txt(s9, Inches(9.3), Inches(2.4), Inches(3.5), Inches(4.0))
para(tf, u'\u2022 Aplikacje r\u00f3\u017cni\u0105 si\u0119 interfejsem',
     Pt(12), False, TEXT_WHITE, space_before=Pt(0))
para(tf, u'\u2022 Odpowiedzi zale\u017c\u0105 od konkretnego rozwi\u0105zania '
     u'(robo-doradca vs aktywny broker)',
     Pt(12), False, TEXT_WHITE, space_before=Pt(8))
para(tf, u'\u2022 Utrudni to znalezienie wsp\u00f3lnych mianownik\u00f3w na etapie analizy',
     Pt(12), False, TEXT_WHITE, space_before=Pt(8))


# ============================================================
# Save the presentation
# ============================================================
output_path = '/projects/sandbox/github/Prezentacja_Fintech.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
